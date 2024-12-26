from langchain_community.document_loaders import PyPDFLoader
from docx import Document as DocxDocument
import csv
from pptx import Presentation
from fastapi import HTTPException

# Função para carregar documentos de diferentes formatos
def load_document(file_path):
    if file_path.endswith('.pdf'):
        loader = PyPDFLoader(file_path)
        return loader.load()
    elif file_path.endswith('.docx'):
        return load_docx(file_path)
    elif file_path.endswith('.txt'):
        return load_txt(file_path)
    elif file_path.endswith('.csv'):
        return load_csv(file_path)
    elif file_path.endswith('.pptx'):
        return load_pptx(file_path)
    else:
        raise HTTPException(status_code=400, detail=f"Formato de arquivo não suportado: {file_path}")

# Carregadores específicos
def load_docx(file_path):
    doc = DocxDocument(file_path)
    return [{"page_content": paragraph.text, "metadata": {"source": file_path}} for paragraph in doc.paragraphs if paragraph.text.strip()]

def load_txt(file_path):
    with open(file_path, "r", encoding="utf-8") as file:
        content = file.read()
    return [{"page_content": content, "metadata": {"source": file_path}}]

def load_csv(file_path):
    rows = []
    with open(file_path, "r", encoding="utf-8") as csvfile:
        reader = csv.reader(csvfile)
        for row in reader:
            rows.append(" ".join(row))
    content = "\n".join(rows)
    return [{"page_content": content, "metadata": {"source": file_path}}]

def load_pptx(file_path):
    prs = Presentation(file_path)
    slides_content = []
    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        slides_content.append(" ".join(slide_text))
    content = "\n".join(slides_content)
    return [{"page_content": content, "metadata": {"source": file_path}}]
