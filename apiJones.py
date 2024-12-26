from fastapi import FastAPI, UploadFile, Form, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from typing import List
import tempfile
import os
import threading
import uuid
import time
import io
import pandas as pd
from docx import Document as DocxDocument
from pptx import Presentation
from langchain.schema import Document
from langchain_core.messages import AIMessage, HumanMessage
from langchain_core.prompts import PromptTemplate
from langchain.chat_models import ChatOpenAI
from langchain.embeddings import OpenAIEmbeddings
from langchain_community.vectorstores import FAISS
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain.chains import create_history_aware_retriever, create_retrieval_chain
from langchain.chains.combine_documents import create_stuff_documents_chain
from langchain_community.document_loaders import PyPDFLoader
from dotenv import load_dotenv
import logging

logging.basicConfig(level=logging.DEBUG)

load_dotenv()

app = FastAPI()

# Adiciona middleware para CORS
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Estrutura para sessões
sessions = {}
lock = threading.Lock()

# Modelo de entrada para perguntas
class QuestionRequest(BaseModel):
    session_id: str
    question: str

# Função para configurar o modelo
def model_openai(model="gpt-4", temperature=0.1):
    llm = ChatOpenAI(model=model, temperature=temperature)
    return llm

# Função para extrair texto dos diferentes formatos de arquivo
def extract_text_from_file(file: UploadFile):
    file_ext = file.filename.split('.')[-1].lower()
    content = file.file.read()

    if file_ext == 'pdf':
        temp_file = tempfile.NamedTemporaryFile(delete=False)
        with open(temp_file.name, 'wb') as f:
            f.write(content)
        loader = PyPDFLoader(temp_file.name)
        return [
            Document(page_content=doc.page_content, metadata={'page': doc.metadata.get('page', 0), 'source': file.filename})
            for doc in loader.load()
        ]

    elif file_ext == 'txt':
        text = content.decode('utf-8')
        return [Document(page_content=text, metadata={'page': 0, 'source': file.filename})]

    elif file_ext == 'csv':
        df = pd.read_csv(io.BytesIO(content))
        text = df.to_string()
        return [Document(page_content=text, metadata={'page': 0, 'source': file.filename})]

    elif file_ext == 'xlsx':
        df = pd.read_excel(io.BytesIO(content))
        text = df.to_string()
        return [Document(page_content=text, metadata={'page': 0, 'source': file.filename})]

    elif file_ext == 'docx':
        doc = DocxDocument(io.BytesIO(content))
        text = '\n'.join([p.text for p in doc.paragraphs if p.text])
        return [Document(page_content=text, metadata={'page': 0, 'source': file.filename})]

    elif file_ext == 'pptx':
        prs = Presentation(io.BytesIO(content))
        docs = []
        for slide_index, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text:
                    slide_text.append(shape.text)
            if slide_text:
                docs.append(Document(
                    page_content='\n'.join(slide_text),
                    metadata={'page': slide_index + 1, 'source': file.filename}
                ))
        return docs

    else:
        raise HTTPException(status_code=400, detail="Formato de arquivo não suportado.")

# Função para configurar o retriever
def config_retriever(uploaded_files):
    docs = []
    for file in uploaded_files:
        docs.extend(extract_text_from_file(file))

    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
    splits = text_splitter.split_documents(docs)

    embeddings = OpenAIEmbeddings()
    vectorstore = FAISS.from_documents(splits, embeddings)
    return vectorstore

# Função para configurar a RAG Chain
def config_rag_chain(retriever):
    llm = model_openai()
    qa_prompt_template = """Você é um assistente virtual prestativo e está respondendo perguntas gerais.
    Use os seguintes pedaços de contexto recuperado para responder à pergunta.
    Se você não sabe a resposta, apenas diga que não sabe. Mantenha a resposta concisa.
    Responda em português. \n\n
    Pergunta: {input} \n
    Contexto: {context}"""

    qa_prompt = PromptTemplate.from_template(qa_prompt_template)
    qa_chain = create_stuff_documents_chain(llm, qa_prompt)

    rag_chain = create_retrieval_chain(retriever, qa_chain)
    return rag_chain

@app.get("/")
async def root():
    return {"message": "Bem-vindo ao sistema de conversa com documentos!"}

@app.post("/upload/")
async def upload_files(files: List[UploadFile]):
    if not files:
        raise HTTPException(status_code=400, detail="Nenhum upload feito.")

    vectorstore = config_retriever(files)
    retriever = vectorstore.as_retriever(search_type="mmr", search_kwargs={"k": 3, "fetch_k": 4})
    rag_chain = config_rag_chain(retriever)
    session_id = str(uuid.uuid4())

    with lock:
        sessions[session_id] = {
            "rag_chain": rag_chain,
            "chat_history": [
                AIMessage(content="Olá, sou o seu assistente virtual! Como posso ajudar você?")
            ]
        }

    return {"session_id": session_id}

@app.post("/ask/")
async def ask_question(request: QuestionRequest):
    session_id = request.session_id
    question = request.question

    with lock:
        session = sessions.get(session_id)
        if not session:
            raise HTTPException(status_code=404, detail="Sessão não encontrada.")

    rag_chain = session["rag_chain"]
    chat_history = session["chat_history"]

    chat_history.append(HumanMessage(content=question))

    start_time = time.time()
    result = rag_chain.invoke({"input": question, "chat_history": chat_history})
    elapsed_time = time.time() - start_time

    response = result['answer']
    context_sources = result['context']

    logging.debug(f"context_sources: {context_sources}")

    # Adiciona a resposta ao histórico
    chat_history.append(AIMessage(content=response))

    return {
        "response": response,
        "chat_history": [{"role": "user", "content": msg.content} if isinstance(msg, HumanMessage)
                         else {"role": "assistant", "content": msg.content} for msg in chat_history],
        "sources": [
            {
                "file": os.path.basename(doc.metadata['source']) if 'source' in doc.metadata else "Fonte desconhecida",
                "page": doc.metadata.get('page', 'Página não especificada')
            }
            for doc in context_sources
        ],
        "response_time": elapsed_time
    }

@app.post("/exit/")
async def exit_chat(request: QuestionRequest):
    session_id = request.session_id

    with lock:
        if session_id in sessions:
            del sessions[session_id]
            return {"message": "Sessão encerrada e contexto apagado."}
        else:
            raise HTTPException(status_code=404, detail="Sessão não encontrada.")
